"""
PowerPoint 자동 생성 유틸리티
python-pptx를 사용한 PPT 자동화
"""

from pathlib import Path
from typing import Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt


class PptGenerator:
    """
    PPT 자동 생성 클래스
    - 기존 템플릿 열기
    - 슬라이드 추가
    - 제목 + 이미지 삽입
    """

    def __init__(self, template_path: Optional[Path] = None):
        """
        Args:
            template_path: 기존 PPT 템플릿 경로 (None이면 빈 PPT)
        """
        if template_path and template_path.exists():
            self.prs = Presentation(str(template_path))
            print(f"[OK] Template loaded: {template_path.name}")
        else:
            self.prs = Presentation()
            print("[OK] New presentation created")

    def add_slide_with_image(
        self,
        title: str,
        image_path: Path,
        layout_index: int = 5  # 5 = Blank layout
    ) -> None:
        """
        제목 + 이미지 슬라이드 추가

        Args:
            title: 슬라이드 제목
            image_path: 삽입할 이미지 경로
            layout_index: 레이아웃 인덱스 (0=Title, 1=Title+Content, 5=Blank)
        """
        # 슬라이드 레이아웃 선택
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목 추가 (상단)
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title

        # 제목 스타일
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.name = 'Arial'

        # 이미지 추가 (제목 아래)
        img_left = Inches(0.5)
        img_top = Inches(1.3)
        img_width = Inches(9)

        if image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                img_left,
                img_top,
                width=img_width
            )
        else:
            print(f"[WARNING] Image not found: {image_path}")

    def save(self, output_path: Path) -> None:
        """
        PPT 파일 저장

        Args:
            output_path: 저장 경로
        """
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        file_size = output_path.stat().st_size / (1024 * 1024)
        print(f"[OK] PPT saved: {output_path.name} ({file_size:.1f} MB)")

    @staticmethod
    def batch_generate_ppt(
        image_files: List[Path],
        output_ppt: Path,
        template_path: Optional[Path] = None,
        title_template: str = "{band} {lna} {port} LNA Gain"
    ) -> None:
        """
        여러 이미지를 하나의 PPT로 자동 생성

        Args:
            image_files: 이미지 파일 리스트
            output_ppt: 출력 PPT 경로
            template_path: PPT 템플릿 (선택)
            title_template: 제목 템플릿
        """
        generator = PptGenerator(template_path)

        for img_path in image_files:
            # 파일명에서 정보 추출 (예: B41_G0_H_ANT1.svg)
            stem = img_path.stem
            parts = stem.split('_')

            if len(parts) >= 3:
                band = parts[0]
                lna = parts[1]
                port = parts[2]
                title = title_template.format(band=band, lna=lna, port=port)
            else:
                title = stem

            print(f"  Adding slide: {title}")
            generator.add_slide_with_image(title, img_path)

        generator.save(output_ppt)
        print(f"[OK] Total slides: {len(image_files)}")
