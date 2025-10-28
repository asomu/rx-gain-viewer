"""
PowerPoint 자동 생성 유틸리티
python-pptx를 사용한 PPT 자동화
"""

from pathlib import Path
from typing import Optional, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PptGenerator:
    """
    PPT 자동 생성 클래스
    - 기존 템플릿 열기
    - 템플릿 레이아웃 자동 감지
    - 슬라이드 추가 (템플릿 스타일 유지)
    """

    def __init__(self, template_path: Optional[Path] = None):
        """
        Args:
            template_path: 기존 PPT 템플릿 경로 (None이면 빈 PPT)
        """
        if template_path and template_path.exists():
            self.prs = Presentation(str(template_path))
            print(f"[OK] Template loaded: {template_path.name}")
            self._detect_best_layout()
        else:
            self.prs = Presentation()
            print("[OK] New presentation created")
            self.layout_index = 5  # Blank layout as default
            self.title_placeholder = None
            self.content_placeholder = None

    def _detect_best_layout(self) -> None:
        """
        템플릿에서 제목+내용 레이아웃 자동 감지
        우선순위: Title and Content > Title Only > Blank
        """
        print("[INFO] Detecting template layouts...")
        
        # 모든 레이아웃 검사
        for idx, layout in enumerate(self.prs.slide_layouts):
            print(f"  Layout {idx}: {layout.name}")
            
            # Title and Content 레이아웃 우선 선택
            if 'title' in layout.name.lower() and 'content' in layout.name.lower():
                self.layout_index = idx
                self._analyze_layout_placeholders(idx)
                print(f"[OK] Selected layout {idx}: {layout.name}")
                return
        
        # Title Only 레이아웃 차선책
        for idx, layout in enumerate(self.prs.slide_layouts):
            if 'title' in layout.name.lower() and 'content' not in layout.name.lower():
                self.layout_index = idx
                self._analyze_layout_placeholders(idx)
                print(f"[OK] Selected layout {idx}: {layout.name} (Title Only)")
                return
        
        # Blank 레이아웃 최후 선택
        self.layout_index = 5 if len(self.prs.slide_layouts) > 5 else 0
        self._analyze_layout_placeholders(self.layout_index)
        print(f"[OK] Selected layout {self.layout_index}: Blank")

    def _analyze_layout_placeholders(self, layout_idx: int) -> None:
        """
        레이아웃의 placeholder 분석
        """
        layout = self.prs.slide_layouts[layout_idx]
        self.title_placeholder = None
        self.content_placeholder = None
        
        for shape in layout.placeholders:
            if shape.placeholder_format.type == 1:  # Title
                self.title_placeholder = shape
                print(f"    Found title placeholder: idx={shape.placeholder_format.idx}")
            elif shape.placeholder_format.type == 7:  # Content
                self.content_placeholder = shape
                print(f"    Found content placeholder: idx={shape.placeholder_format.idx}")

    def add_slide_with_image(
        self,
        title: str,
        image_path: Path,
        layout_index: Optional[int] = None
    ) -> None:
        """
        제목 + 이미지 슬라이드 추가 (템플릿 레이아웃 사용)

        Args:
            title: 슬라이드 제목
            image_path: 삽입할 이미지 경로
            layout_index: 레이아웃 인덱스 (None이면 자동 감지된 레이아웃 사용)
        """
        # 레이아웃 선택
        layout_idx = layout_index if layout_index is not None else self.layout_index
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목 설정 (placeholder 사용 또는 수동 추가)
        title_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Title placeholder
                title_shape = shape
                title_shape.text = title
                break
        
        # Title placeholder가 없으면 수동으로 추가
        if title_shape is None:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.8)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            title_shape.text_frame.text = title
            title_shape.text_frame.paragraphs[0].font.size = Pt(28)
            title_shape.text_frame.paragraphs[0].font.bold = True

        # 이미지 추가
        if not image_path.exists():
            print(f"[WARNING] Image not found: {image_path}")
            return

        # Content placeholder가 있으면 그 위치에 이미지 삽입
        content_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 7:  # Content placeholder
                content_shape = shape
                break

        if content_shape:
            # Content placeholder 위치 사용
            left = content_shape.left
            top = content_shape.top
            width = content_shape.width
            
            # Placeholder 삭제하고 이미지로 대체
            sp = content_shape.element
            sp.getparent().remove(sp)
            
            slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=width
            )
        else:
            # Content placeholder가 없으면 제목 아래에 이미지 배치
            img_left = Inches(0.5)
            img_top = Inches(1.3)
            img_width = Inches(9)
            
            slide.shapes.add_picture(
                str(image_path),
                img_left,
                img_top,
                width=img_width
            )

    def save(self, output_path: Path) -> None:
        """
        PPT 파일 저장

        Args:
            output_path: 저장 경로
        """
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        file_size = output_path.stat().st_size / (1024 * 1024)
        print(f"[OK] PPT saved: {output_path.name} ({file_size:.1f} MB)")

    @staticmethod
    def batch_generate_ppt(
        image_files: List[Path],
        output_ppt: Path,
        template_path: Optional[Path] = None,
        title_template: str = "{band} {lna} {port} LNA Gain"
    ) -> None:
        """
        여러 이미지를 하나의 PPT로 자동 생성

        Args:
            image_files: 이미지 파일 리스트
            output_ppt: 출력 PPT 경로
            template_path: PPT 템플릿 (선택)
            title_template: 제목 템플릿
        """
        generator = PptGenerator(template_path)

        for img_path in image_files:
            # 파일명에서 정보 추출 (예: B41_G0_H_ANT1.png)
            stem = img_path.stem
            parts = stem.split('_')

            if len(parts) >= 3:
                band = parts[0]
                lna = parts[1]
                port = parts[2]
                title = title_template.format(band=band, lna=lna, port=port)
            else:
                title = stem

            print(f"  Adding slide: {title}")
            generator.add_slide_with_image(title, img_path)

        generator.save(output_ppt)
        print(f"[OK] Total slides: {len(image_files)}")
